from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Create Presentation
prs = Presentation()

# Brand Colors
PRIMARY_COLOR = RGBColor(15, 44, 89)    # #0f2c59 (Navy)
SECONDARY_COLOR = RGBColor(209, 166, 98) # #d1a662 (Gold)
ACCENT_COLOR = RGBColor(248, 249, 250)   # #f8f9fa (Light Gray)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(50, 50, 50)

def clear_slide(slide):
    """Remove default placeholders to have a clean slate"""
    for shape in slide.placeholders:
        sp = shape.element
        sp.getparent().remove(sp)

def add_header_footer(slide, title_text, page_num=None):
    # Header Line
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.05)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + Inches(0.8), width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = SECONDARY_COLOR
    shape.line.fill.background() # No line

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = 'Arial'
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    # Footer
    if page_num:
        footer_box = slide.shapes.add_textbox(Inches(8.5), Inches(7.2), Inches(1), Inches(0.3))
        tf = footer_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(page_num)
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(150, 150, 150)
        p.alignment = PP_ALIGN.RIGHT
    
    # Logo mark (simulated text)
    logo_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.2), Inches(2), Inches(0.3))
    tf = logo_box.text_frame
    p = tf.paragraphs[0]
    p.text = "FATIC Presentation"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(150, 150, 150)

def add_image_with_style(slide, image_path, left, top, height):
    if os.path.exists(image_path):
        # Shadow effect (Rectangle behind)
        shadow_left = left + Pt(5)
        shadow_top = top + Pt(5)
        shadow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, shadow_left, shadow_top, height * 1.6, height) # Approx 16:9 ratio
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = SECONDARY_COLOR
        shadow.line.fill.background()
        
        # Image
        slide.shapes.add_picture(image_path, left, top, height=height)

# --- SLIDE 1: TITLE SLIDE ---
slide_layout = prs.slide_layouts[6] # Blank
slide = prs.slides.add_slide(slide_layout)

# Background
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5))
bg.fill.solid()
bg.fill.fore_color.rgb = PRIMARY_COLOR
bg.line.fill.background()

# Decorative Circle
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6), Inches(-2), Inches(6), Inches(6))
circle.fill.solid()
circle.fill.fore_color.rgb = SECONDARY_COLOR
circle.fill.transparency = 0.8
circle.line.fill.background()

# Title Content
title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "FATIC"
p.font.name = 'Arial Black'
p.font.size = Pt(66)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.LEFT

p2 = tf.add_paragraph()
p2.text = "KIẾN TẠO GIÁ TRỊ BỀN VỮNG"
p2.font.name = 'Arial'
p2.font.size = Pt(24)
p2.font.bold = True
p2.font.color.rgb = SECONDARY_COLOR
p2.space_before = Pt(10)

p3 = tf.add_paragraph()
p3.text = "Giải Pháp Quản Lý Vận Hành & Bảo Trì Kỹ Thuật"
p3.font.name = 'Arial'
p3.font.size = Pt(18)
p3.font.color.rgb = RGBColor(230, 230, 230)
p3.space_before = Pt(20)

# --- SLIDE 2: ABOUT US ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_footer(slide, "Giới Thiệu Chung", 2)

# Content Box
textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.5), Inches(5))
tf = textbox.text_frame
tf.word_wrap = True

points = [
    ("Lịch Sử", "Thành lập 2016 - Hơn 10 năm phát triển."),
    ("Vị Thế", "Đơn vị Top 10 về Quản lý Vận hành & Bảo trì."),
    ("Quy Mô", "500+ Nhân sự chuyên gia & Kỹ thuật viên cao cấp."),
    ("Cam Kết", "Chuyên nghiệp - Minh bạch - Hiệu quả.")
]

for title, desc in points:
    p = tf.add_paragraph()
    p.text = title
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = PRIMARY_COLOR
    p.space_before = Pt(12)
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(14)
    p2.font.color.rgb = DARK_GRAY
    p2.space_after = Pt(8)

# Image
add_image_with_style(slide, r"e:\code\fatic\ppt_images\home.png", Inches(5.2), Inches(2), Inches(3.5))

# --- SLIDE 3: VISION & MISSION ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_footer(slide, "Tầm Nhìn & Sứ Mệnh", 3)

# 2 Column Layout
# Left: Vision
left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.5), Inches(2))
tf = left_box.text_frame
p = tf.paragraphs[0]
p.text = "TẦM NHÌN"
p.font.bold = True
p.font.size = Pt(20)
p.font.color.rgb = SECONDARY_COLOR
p.space_after = Pt(10)

p2 = tf.add_paragraph()
p2.text = "Trở thành biểu tượng uy tín hàng đầu Việt Nam trong lĩnh vực cung cấp giải pháp quản lý Bất động sản toàn diện."
p2.font.size = Pt(16)
p2.font.color.rgb = DARK_GRAY

# Right: Mission
right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4.3), Inches(2))
tf = right_box.text_frame
p = tf.paragraphs[0]
p.text = "SỨ MỆNH"
p.font.bold = True
p.font.size = Pt(20)
p.font.color.rgb = SECONDARY_COLOR
p.space_after = Pt(10)

missions = [
    "Dịch vụ chuẩn mực quốc tế.",
    "Tối ưu chi phí vận hành cho CĐT.",
    "Kiến tạo cộng đồng văn minh."
]
for m in missions:
    pm = tf.add_paragraph()
    pm.text = "• " + m
    pm.font.size = Pt(16)
    pm.font.color.rgb = DARK_GRAY

# Bottom: Core Values
val_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
tf = val_box.text_frame
p = tf.paragraphs[0]
p.text = "GIÁ TRỊ CỐT LÕI"
p.font.bold = True
p.font.size = Pt(20)
p.font.color.rgb = PRIMARY_COLOR
p.alignment = PP_ALIGN.CENTER
p.space_after = Pt(20)

# Value shapes
values = ["CHUYÊN NGHIỆP", "TRÁCH NHIỆM", "TẬN TÂM", "HIỆU QUẢ"]
start_x = 0.5
width_x = 2
gap = 0.3
for val in values:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(start_x), Inches(5.5), Inches(width_x), Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
    shape.line.color.rgb = SECONDARY_COLOR
    
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = val
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.alignment = PP_ALIGN.CENTER
    start_x += width_x + gap

# --- SLIDE 4: SERVICES ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_footer(slide, "Dịch Vụ Cốt Lõi", 4)

textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.5), Inches(5))
tf = textbox.text_frame
services = [
    ("Quản Lý Tòa Nhà", "Vận hành chung cư, văn phòng, TTTM theo chuẩn ISO."),
    ("Bảo Trì Kỹ Thuật", "Hệ thống M&E, PCCC, HVAC, Thang máy."),
    ("Tư Vấn Giải Pháp", "Setup quy trình, hoạch định tài chính tối ưu."),
    ("Cung Ứng Nhân Sự", "Bảo vệ, Vệ sinh, Kỹ thuật viên lành nghề.")
]

for title, desc in services:
    p = tf.add_paragraph()
    p.text = "➤ " + title
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = PRIMARY_COLOR
    p.space_after = Pt(2)
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(13)
    p2.font.color.rgb = DARK_GRAY
    p2.level = 0
    p2.space_after = Pt(12)

add_image_with_style(slide, r"e:\code\fatic\ppt_images\services.png", Inches(5.2), Inches(2), Inches(3.5))

# --- SLIDE 5: PROJECTS ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_footer(slide, "Dự Án & Đối Tác", 5)

textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.5), Inches(5))
tf = textbox.text_frame
partners = [
    "Đối Tác Chiến Lược:",
    "VinHomes, NovaLand, Sun Group, Bitexco, Ecopark.",
    "",
    "Lĩnh Vực Triển Khai:",
    "• Khu đô thị & Chung cư cao cấp.",
    "• Tòa nhà văn phòng hạng A.",
    "• Khu công nghiệp công nghệ cao."
]

for line in partners:
    p = tf.add_paragraph()
    p.text = line
    if "Đối Tác" in line or "Lĩnh Vực" in line:
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = PRIMARY_COLOR
        p.space_before = Pt(10)
    else:
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY

add_image_with_style(slide, r"e:\code\fatic\ppt_images\projects.png", Inches(5.2), Inches(2), Inches(3.5))

# --- SLIDE 6: CONTACT ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
# Header
add_header_footer(slide, "Liên Hệ Hợp Tác", 6)

# Contact Box
bg_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2), Inches(4.5), Inches(3.5))
bg_box.fill.solid()
bg_box.fill.fore_color.rgb = RGBColor(245, 245, 245)
bg_box.line.color.rgb = RGBColor(200, 200, 200)

tf = bg_box.text_frame
tf.margin_left = Inches(0.2)
tf.margin_top = Inches(0.2)

contacts = [
    ("Trụ Sở Chính (Hà Nội)", "88 Láng Hạ, Đống Đa, Hà Nội"),
    ("Hotline", "(+84) 24 3999 8888"),
    ("Chi Nhánh TP.HCM", "Bitexco Tower, Quận 1"),
    ("Hotline", "(+84) 28 3999 9999"),
    ("Email / Website", "contact@fatic.com.vn\nwww.fatic.com.vn")
]

for label, info in contacts:
    p = tf.add_paragraph()
    p.text = label
    p.font.bold = True
    p.font.size = Pt(12)
    p.font.color.rgb = PRIMARY_COLOR
    
    p2 = tf.add_paragraph()
    p2.text = info
    p2.font.size = Pt(12)
    p2.font.color.rgb = DARK_GRAY
    p2.space_after = Pt(8)

add_image_with_style(slide, r"e:\code\fatic\ppt_images\contact.png", Inches(5.2), Inches(2), Inches(3.5))

# --- SAVE ---
output_file = 'Gioi_Thieu_FATIC_V2.pptx'
prs.save(output_file)
print(f"Successfully created {output_file}")
